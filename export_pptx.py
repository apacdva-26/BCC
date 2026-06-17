"""
PPTX export — SAP Business Case Calculator
Slide 1: Value Tree  (one column per selected business area)
Slide 2: Project Economics (chart + KPI boxes)
"""
from __future__ import annotations
import copy, io
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

TEMPLATE_PATH = Path(__file__).parent / "Presentation1.pptx"

# ── Colors ────────────────────────────────────────────────────────────────────
C_BLUE   = RGBColor(0x00, 0x70, 0xF2)
C_NAVY   = RGBColor(0x00, 0x14, 0x4A)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_REC    = RGBColor(0xD1, 0xEF, 0xFF)   # recurring  – light blue
C_ONE    = RGBColor(0xE2, 0xD8, 0xFF)   # one-time   – light purple

EMU_PER_IN = 914400


# ── Helpers ───────────────────────────────────────────────────────────────────

def _rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def _fmt_m(value: float, symbol: str) -> str:
    m = value / 1_000_000
    if abs(m) >= 1000:
        return f"{symbol}{m/1000:,.1f}B"
    elif abs(m) >= 1:
        return f"{symbol}{m:,.1f}M"
    elif abs(value) >= 1000:
        return f"{symbol}{value/1000:,.1f}K"
    else:
        return f"{symbol}{value:,.0f}"


def _set_pos(shape, left_in: float, top_in: float,
             w_in: float | None = None, h_in: float | None = None):
    sp = shape._element
    xfrm = sp.find('.//' + qn('p:xfrm')) or sp.find('.//' + qn('a:xfrm'))
    if xfrm is None:
        return
    off = xfrm.find(qn('a:off'))
    ext = xfrm.find(qn('a:ext'))
    if off is not None:
        off.set('x', str(int(left_in * EMU_PER_IN)))
        off.set('y', str(int(top_in  * EMU_PER_IN)))
    if ext is not None:
        if w_in is not None:
            ext.set('cx', str(int(w_in * EMU_PER_IN)))
        if h_in is not None:
            ext.set('cy', str(int(h_in * EMU_PER_IN)))


def _clear_set(shape, lines: list[tuple[str, float, bool, RGBColor]],
               align=PP_ALIGN.LEFT, font_name: str = "72 Brand"):
    """Replace all text. lines = [(text, pt, bold, color)]"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, (text, pt, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color


def _clone(slide, shape):
    sp_tree = slide.shapes._spTree
    new_el = copy.deepcopy(shape._element)
    sp_tree.append(new_el)
    return slide.shapes[-1]


def _set_cell(cell, text: str, pt: float, bold: bool,
              color: RGBColor, align=PP_ALIGN.LEFT,
              fill: RGBColor | None = None,
              font_name: str = "72 Brand"):
    tf = cell.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    if text:
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color
    if fill is not None:
        tc = cell._tc
        tcPr = tc.find(qn('a:tcPr'))
        if tcPr is None:
            tcPr = etree.SubElement(tc, qn('a:tcPr'))
        for ex in tcPr.findall(qn('a:solidFill')):
            tcPr.remove(ex)
        sf  = etree.SubElement(tcPr, qn('a:solidFill'))
        clr = etree.SubElement(sf, qn('a:srgbClr'))
        clr.set('val', _rgb_hex(fill))


def _resize_table(tbl_shape, n_rows: int, row_h_in: float = 0.32):
    table  = tbl_shape.table
    tbl_el = table._tbl
    target = int(row_h_in * EMU_PER_IN)
    existing = list(tbl_el.findall(qn('a:tr')))
    cur = len(existing)
    if n_rows > cur:
        ref = existing[-1]
        for _ in range(n_rows - cur):
            tbl_el.append(copy.deepcopy(ref))
    elif n_rows < cur:
        for row_el in existing[n_rows:]:
            tbl_el.remove(row_el)
    for tr in tbl_el.findall(qn('a:tr')):
        tr.set('h', str(target))
    xfrm = tbl_shape._element.find('.//' + qn('p:xfrm')) or \
           tbl_shape._element.find('.//' + qn('a:xfrm'))
    if xfrm is not None:
        ext = xfrm.find(qn('a:ext'))
        if ext is not None:
            ext.set('cy', str(target * n_rows))


# ── Slide 1: Value Tree ───────────────────────────────────────────────────────

def _build_slide1(prs: Presentation, results: dict,
                  customer_name: str, currency_symbol: str, currency_code: str,
                  industry: str, maturity_label: str):

    slide = prs.slides[0]
    sp_tree = slide.shapes._spTree

    # Grab template shapes (first instance of each repeated shape)
    def _first(name):
        return next(s for s in slide.shapes if s.name == name)
    def _all(name):
        return [s for s in slide.shapes if s.name == name]

    tmpl_lob   = _all('Lob_Name_1')[0]
    tmpl_ob1   = _all('ob1')[0]
    tmpl_rb1   = _all('rb1')[0]
    tmpl_tbl   = _all('tbl_lob2')[0]

    # Template block reference geometry
    LOB_W  = tmpl_lob.width  / EMU_PER_IN
    LOB_H  = tmpl_lob.height / EMU_PER_IN
    OB1_W  = tmpl_ob1.width  / EMU_PER_IN
    OB1_H  = tmpl_ob1.height / EMU_PER_IN
    TBL_W  = tmpl_tbl.width  / EMU_PER_IN

    # Remove all existing LOB blocks (keep structural shapes)
    KEEP = {
        'Round Same Side Corner Rectangle 48',
        'Rectangle: Top Corners Rounded 60',
        'Rectangle: Top Corners Rounded 61',
        'Rectangle 62', 'Group 63', 'Title 1',
        'Left Bracket 98',
        'Lob_Name_4',
        'Rectangle: Rounded Corners 99',
        'Rectangle: Rounded Corners 100',
        'TextBox 24',
        'Rounded Rectangle 7', 'Rounded Rectangle 8',
        'Rounded Rectangle 9', 'Rounded Rectangle 10',
    }
    for s in [s for s in slide.shapes if s.name not in KEEP]:
        sp_tree.remove(s._element)

    # Update slide title
    _clear_set(_first('Title 1'),
               [("Value Tree", 14, True, RGBColor(0x00, 0x2A, 0x86))],
               align=PP_ALIGN.LEFT, font_name="72 Brand Medium")

    # Update info sidebar
    _clear_set(_first('Rectangle 62'), [
        (f"Currency: {currency_code}",  7, True,  C_WHITE),
        (f"Industry: {industry}",        7, False, C_WHITE),
        (f"Maturity: {maturity_label}",  7, False, C_WHITE),
        ("*Illustrative estimates based on industry benchmarks.", 6, False, C_WHITE),
    ], font_name="72 Brand")

    # Determine active areas & layout
    AREA_ORDER = ["Asset Management", "Finance", "Manufacturing",
                  "Sales", "Procurement", "R&D", "Supply Chain"]
    active = [a for a in AREA_ORDER if a in results["by_area"]]
    n = len(active)

    # Slide usable content area (left of the info sidebar)
    SLIDE_W    = 10.70   # inches available for LoB columns
    START_LEFT = 0.15
    CONTENT_TOP = 2.40   # below total summary row
    CONTENT_BOT = 7.20
    CONTENT_H   = CONTENT_BOT - CONTENT_TOP

    # Layout: 1 or 2 rows
    if n <= 4:
        rows_layout = [active]
    else:
        split = (n + 1) // 2
        rows_layout = [active[:split], active[split:]]

    n_rows_layout = len(rows_layout)
    ROW_GAP  = 0.12
    ROW_H    = (CONTENT_H - ROW_GAP * (n_rows_layout - 1)) / n_rows_layout

    # Header label top & metric box geometry
    METRIC_GAP = 0.06
    METRIC_H   = OB1_H
    TBL_ROW_H  = 0.30

    def _place_block(area: str, left: float, top: float, col_w: float):
        data = results["by_area"].get(area)
        if not data:
            return

        rec     = data["subtotal_recurring"]
        one     = data["subtotal_one_time"]
        drivers = data["drivers"]

        # LoB name label (centered in column)
        lob = _clone(slide, tmpl_lob)
        _set_pos(lob, left, top, col_w, LOB_H)
        _clear_set(lob, [(area, 9, True, C_BLACK)],
                   align=PP_ALIGN.CENTER, font_name="72 Brand Medium")

        metric_top = top + LOB_H + 0.04
        metric_w   = (col_w - METRIC_GAP * 3) / 2

        # Recurring box (ob1) — light blue
        ob = _clone(slide, tmpl_ob1)
        _set_pos(ob, left + METRIC_GAP, metric_top, metric_w, METRIC_H)
        _clear_set(ob, [(_fmt_m(rec, currency_symbol), 9, True, C_BLACK)],
                   align=PP_ALIGN.CENTER, font_name="72 Brand")

        # One-time box (rb1) — light purple
        rb = _clone(slide, tmpl_rb1)
        _set_pos(rb, left + METRIC_GAP * 2 + metric_w, metric_top, metric_w, METRIC_H)
        _clear_set(rb,
                   [(_fmt_m(one, currency_symbol) if one > 0 else "—", 9, True, C_BLACK)],
                   align=PP_ALIGN.CENTER, font_name="72 Brand")

        # Detail table
        tbl_top  = metric_top + METRIC_H + 0.06
        tbl_h_avail = (top + ROW_H) - tbl_top - 0.05
        n_drv    = max(len(drivers), 1)
        row_h    = min(TBL_ROW_H, tbl_h_avail / n_drv)

        tbl = _clone(slide, tmpl_tbl)
        _set_pos(tbl, left + METRIC_GAP * 0.5, tbl_top, col_w - METRIC_GAP, None)
        _resize_table(tbl, n_drv, row_h)

        col_widths_in = [col_w * 0.55, col_w * 0.18, col_w * 0.27]
        tbl_el = tbl.table._tbl
        tblGrid = tbl_el.find(qn('a:tblGrid'))
        if tblGrid is not None:
            cols = tblGrid.findall(qn('a:gridCol'))
            for ci, cw in enumerate(col_widths_in):
                if ci < len(cols):
                    cols[ci].set('w', str(int(cw * EMU_PER_IN)))

        for ri, d in enumerate(drivers):
            fill_color = C_ONE if d["one_time"] else C_REC
            row_obj = tbl.table.rows[ri]
            _set_cell(row_obj.cells[0], d["name"],
                      7, False, C_BLACK, PP_ALIGN.LEFT, font_name="72 Brand")
            _set_cell(row_obj.cells[1], f"{d['improvement_pct']:.0f}%",
                      7, False, C_BLACK, PP_ALIGN.CENTER, font_name="72 Brand")
            _set_cell(row_obj.cells[2], _fmt_m(d["benefit"], currency_symbol),
                      7, True, C_BLACK, PP_ALIGN.CENTER, fill_color, font_name="72 Brand")

    # Place all blocks
    for ri, row_areas in enumerate(rows_layout):
        nc     = len(row_areas)
        col_w  = SLIDE_W / nc
        row_top = CONTENT_TOP + ri * (ROW_H + ROW_GAP)
        for ci, area in enumerate(row_areas):
            _place_block(area, START_LEFT + ci * col_w, row_top, col_w)

    # Total summary row
    total_rec = results["annual_recurring_benefit"]
    total_one = results["one_time_benefit"]
    prod_ben  = results["productivity"]["annual_benefit"]
    total_all = total_rec + total_one

    lob4 = _first('Lob_Name_4')
    _set_pos(lob4, START_LEFT, 1.44, SLIDE_W, 0.50)
    _clear_set(lob4, [
        (f"Total Steady-state Value Potential for {customer_name}", 12, True, C_BLUE),
    ], align=PP_ALIGN.LEFT, font_name="72 Brand Medium")

    r99 = _first('Rectangle: Rounded Corners 99')
    _set_pos(r99, START_LEFT + SLIDE_W * 0.42, 1.98, 2.10, 0.34)
    _clear_set(r99,
               [(f"Recurring: {_fmt_m(total_rec + prod_ben, currency_symbol)}", 10, True, C_BLACK)],
               align=PP_ALIGN.CENTER, font_name="72 Brand")

    r100 = _first('Rectangle: Rounded Corners 100')
    _set_pos(r100, START_LEFT + SLIDE_W * 0.42 + 2.20, 1.98, 2.10, 0.34)
    _clear_set(r100,
               [(f"One-time: {_fmt_m(total_one, currency_symbol)}", 10, True, C_BLACK)],
               align=PP_ALIGN.CENTER, font_name="72 Brand")


# ── Slide 2: Project Economics ────────────────────────────────────────────────

def _build_slide2(prs: Presentation, results: dict,
                  currency_symbol: str, currency_code: str,
                  realization_recurring: list[float],
                  realization_onetime: list[float],
                  discount_rate_pct: float):

    from pptx.chart.data import ChartData

    slide = prs.slides[1]

    acv_by_year     = results["acv_by_year"]
    impl_cost       = results["impl_cost"]
    rec_annual      = results["annual_recurring_benefit"]
    one_total       = results["one_time_benefit"]
    real_rec        = [r / 100.0 for r in realization_recurring]
    real_one        = [r / 100.0 for r in realization_onetime]

    rec_by_year     = [rec_annual * real_rec[i] for i in range(5)]
    one_by_year     = [one_total  * real_one[i] for i in range(5)]
    cost_by_year    = [acv_by_year[i] + (impl_cost if i == 0 else 0) for i in range(5)]
    benefit_by_year = results["benefit_by_year"]

    cumulative = []
    cum = 0.0
    for i in range(5):
        cum += benefit_by_year[i] - cost_by_year[i]
        cumulative.append(cum)

    # Chart data
    chart_shape = next(s for s in slide.shapes if s.name == 'Chart 17')
    cd = ChartData()
    cd.categories = ['Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5']
    cd.add_series('SAP Subscription Cost',          [-v / 1e6 for v in acv_by_year])
    cd.add_series('Partner Implementation Cost',    [-impl_cost / 1e6, 0, 0, 0, 0])
    cd.add_series('Recurring Annual Benefits',       [v / 1e6 for v in rec_by_year])
    cd.add_series('One-time Free Cashflow Benefits', [v / 1e6 for v in one_by_year])
    cd.add_series('Cumulative Benefits',             [v / 1e6 for v in cumulative])
    chart_shape.chart.replace_data(cd)

    # KPI boxes
    npv         = results["npv"]
    npv_roi_pct = results["npv_roi_pct"]
    payback     = results["payback_years"]
    irr         = results.get("irr_pct")
    irr_text    = f"{irr:.0f}%" if irr is not None else "N/A"

    def _set_kpi(shape_name: str, val_text: str, label_text: str):
        shape = next(s for s in slide.shapes if s.name == shape_name)
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = val_text
        r1.font.name = "72 Brand"
        r1.font.size = Pt(16)
        r1.font.bold = True
        r1.font.color.rgb = C_BLUE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label_text
        r2.font.name = "72 Brand"
        r2.font.size = Pt(11)
        r2.font.bold = True
        r2.font.color.rgb = C_BLACK

    _set_kpi('Rectangle: Rounded Corners 20', _fmt_m(npv, currency_symbol), '5-year NPV')
    _set_kpi('Rectangle: Rounded Corners 21', f"{npv_roi_pct:.0f}%",         '5-year ROI')
    _set_kpi('Rectangle: Rounded Corners 22', f"{payback:.1f} Yrs",          'Payback Period')

    # Breakeven label
    be = next(s for s in slide.shapes if s.name == 'TextBox 18')
    be.text_frame.paragraphs[0].runs[0].text = f"Breakeven at {payback:.1f} years"

    # Assumptions box
    rec_str = " / ".join(f"{r:.0f}%Y{i+1}" for i, r in enumerate(realization_recurring))
    one_str = " / ".join(f"{r:.0f}%Y{i+1}" for i, r in enumerate(realization_onetime))
    assump = next(s for s in slide.shapes if s.name == 'Rectangle 23')
    _clear_set(assump, [
        ("Assumptions for Benefits:",                                           7,   True,  C_BLACK),
        ("Benefits estimated based on outside-in, conservative benchmarks.",    6.5, False, C_BLACK),
        ("Annual amounts shown in nominal value (no inflation adjustment).",    6.5, False, C_BLACK),
        (f"Recurring realization: {rec_str}",                                   6.5, False, C_BLACK),
        (f"One-time realization:  {one_str}",                                   6.5, False, C_BLACK),
        ("",                                                                    4,   False, C_BLACK),
        ("Assumptions for Costs:",                                              7,   True,  C_BLACK),
        ("All costs are illustrative only.",                                    6.5, False, C_BLACK),
        (f"Discount rate: {discount_rate_pct:.0f}% for NPV calculation.",       6.5, False, C_BLACK),
        ("Subscription costs as per TCO comparison.",                           6.5, False, C_BLACK),
        ("*Illustrative estimates only. Further discovery required.",           6,   False, C_BLACK),
    ], font_name="72 Brand")


# ── Public API ────────────────────────────────────────────────────────────────

def generate_pptx(
    results: dict,
    customer_name: str,
    currency_symbol: str,
    currency_code: str,
    industry: str,
    maturity_label: str,
    realization_recurring: list[float],
    realization_onetime: list[float],
    discount_rate_pct: float,
) -> bytes:
    prs = Presentation(str(TEMPLATE_PATH))
    _build_slide1(prs, results, customer_name, currency_symbol, currency_code,
                  industry, maturity_label)
    _build_slide2(prs, results, currency_symbol, currency_code,
                  realization_recurring, realization_onetime, discount_rate_pct)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
